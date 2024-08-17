from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.middleware.cors import CORSMiddleware
import fitz   # PyMuPDF
from pptx import Presentation
from transformers import pipeline

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize summarizer pipeline
# try:
# except RuntimeError as e:
#     raise RuntimeError("Ensure you have PyTorch or TensorFlow installed") from e
summarizer = pipeline("summarization")

@app.post("/upload")
async def upload(file: UploadFile = File(...)):
    if file.content_type == "application/pdf":
        text = extract_text_from_pdf(file)
    elif file.content_type == "text/plain":
        text = (await file.read()).decode('utf-8')
    elif file.content_type in ["application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
        text = extract_text_from_ppt(file)
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type")

    summary = summarize_text(text)
    return {"summary": summary}

def extract_text_from_pdf(file):
    doc = fitz.open(stream=file.file.read(), filetype='pdf')
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def extract_text_from_ppt(file):
    prs = Presentation(file.file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

def summarize_text(text):
    max_chunk_size = 1000  # Adjust based on your model's limits
    chunks = [text[i:i+max_chunk_size] for i in range(0, len(text), max_chunk_size)]
    summaries = [summarizer(chunk)[0]['summary_text'] for chunk in chunks]
    return ' '.join(summaries)

if __name__ == '__main__':
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
