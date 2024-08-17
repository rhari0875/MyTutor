from fastapi import FastAPI, File, UploadFile, HTTPException, Request, Form
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware
import os
from fastapi.responses import FileResponse
import pypandoc
from docx import Document
from pypdf import PdfReader
from pptx import Presentation
from fpdf import FPDF
import random
import string
import json

# from moviepy.editor import VideoFileClip
# from quiz_generator import export_quiz
# from flash_card_generator import export_flashcards
# from summary_generator import export_summary
from Summary.export_summary import export_summary
from Quiz.export_quiz import export_quiz
from FlashCards.export_flashcards import export_flashcards
from gemini import prompt_everyting
# from speech_to_text import get_audio

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# def generate_quiz(flashcards):
#     quiz = []
#     for card in flashcards:
#         if random.choice([True, False]):
#             question = {"question": card[1], "possible_answers": [], "index": -1}
#             incorrect_answers = [other_card[0] for other_card in flashcards if other_card != card]
#             question["possible_answers"] = random.sample(incorrect_answers, 3)
#             question["index"] = random.randint(0, 3)
#             question["possible_answers"].insert(question["index"], card[0])
#         else:
#             question = {"question": card[0], "possible_answers": [], "index": -1}
#             incorrect_answers = [other_card[1] for other_card in flashcards if other_card != card]
#             question["possible_answers"] = random.sample(incorrect_answers, 3)
#             question["index"] = random.randint(0, 3)
#             question["possible_answers"].insert(question["index"], card[1])
#         quiz.append(question)
#     return quiz

# def generate_id():
#     ids = set()
#     for file in os.listdir(os.getcwd() + "/database"):
#         if file.endswith(".json"):
#             ids.add(file)

#     letters = string.ascii_letters + string.digits
#     id = "".join(random.choice(letters) for _ in range(5))
#     while id in ids:
#         id = "".join(random.choice(letters) for _ in range(5))
#     return id

# def example_response():
#     flash_cards = [
#         ["Photosynthesis", "The process by which green plants and some other organisms use sunlight to synthesize foods with the help of chlorophyll."],
#         ["Mitochondria", "Organelles that generate most of the chemical energy needed to power the biochemical reactions of cells."],
#         ["Newton's Laws", "Three fundamental principles of classical mechanics proposed by Sir Isaac Newton."],
#         ["H2O", "Chemical formula for water, composed of two hydrogen atoms bonded to one oxygen atom."],
#         ["Civil Rights Movement", "A struggle for social justice that took place mainly during the 1950s and 1960s for African Americans to gain equal rights under the law in the United States."],
#         ["Palindrome", "A word, phrase, number, or other sequence of characters that reads the same forward and backward."],
#         ["The Great Depression", "A severe worldwide economic depression that took place mostly during the 1930s."],
#         ["E=mc^2", "Albert Einstein's famous equation, which expresses the relationship between energy (E), mass (m), and the speed of light (c) squared."],
#         ["Renaissance", "A period in European history marking the transition from the Middle Ages to modernity and covering the 14th to 17th centuries."],
#         ["Cell Division", "The process by which a parent cell divides into two or more daughter cells."],
#     ]
    # summary = """Lorem ipsum dolor sit amet, consectetur adipiscing elit, sed do eiusmod tempor incididunt ut labore et dolore magna aliqua. Pharetra vel turpis nunc eget lorem dolor sed. Diam vulputate ut pharetra sit amet. Amet dictum sit amet justo. Nibh praesent tristique magna sit amet purus gravida quis. Nunc id cursus metus aliquam eleifend mi in nulla posuere. Sollicitudin aliquam ultrices sagittis orci. Egestas erat imperdiet sed euismod nisi porta lorem mollis. In fermentum posuere urna nec tincidunt praesent semper feugiat nibh. Vel pharetra vel turpis nunc eget lorem. Suspendisse sed nisi lacus sed viverra tellus in hac habitasse. Adipiscing enim eu turpis egestas. Facilisis gravida neque convallis a cras semper. Quam quisque id diam vel quam elementum pulvinar etiam."""

    # quiz = generate_quiz(flash_cards)

    # id = generate_id()
    # letters = string.ascii_letters + string.digits
    # title = "Title " + "".join(random.choice(letters) for _ in range(8))
    # db_json = {
    #     "summary": summary,
    #     "flash_cards": flash_cards,
    #     "quiz": quiz,
    #     "title": title,
    # }

    # with open(os.getcwd() + "/database/" + id + ".json", "w") as f:
    #     json.dump(db_json, f)

    # return {"id": id}

def handle_pdf(file_path):
    reader = PdfReader(file_path)
    number_of_pages = len(reader.pages)
    s = ""
    for i in range(number_of_pages):
        page = reader.pages[i]
        text = page.extract_text()
        s += text
    return s

def handle_txt(file_path):
    with open(file_path, "r") as f:
        s = f.read()
    return s

def handle_docx(file_path):
    d = Document(file_path)
    s = ""
    for paragraph in d.paragraphs:
        s += paragraph.text + "\n"
    return s

# def handle_mp3(file_path):
#     s = get_audio(file_path)
#     return s

# def handle_mp4(file_path):
#     video = VideoFileClip(file_path)
#     audio = video.audio
#     audio_file_path = "file.mp3"
#     audio.write_audiofile(audio_file_path)
#     audio.close()
#     video.close()
#     return handle_mp3(audio_file_path)

def handle_pptx(file_path):
    p = Presentation(file_path)
    s = ""
    for slide in p.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                s += shape.text + "\n"
    return s

@app.get("/")
async def hello_world():
    return {"message": "Hello, World!"}

@app.post("/upload")
async def upload(file: UploadFile = File(...)):
    name = file.filename
    extension = name.split(".")[-1]
    file_path = f"file.{extension}"
    
    with open(file_path, "wb") as buffer:
        buffer.write(file.file.read())

    if extension == "pdf":
        s = handle_pdf(file_path)
    elif extension == "txt":
        s = handle_txt(file_path)
    elif extension == "docx":
        s = handle_docx(file_path)
    # elif extension == "mp3":
    #     s = handle_mp3(file_path)
    # elif extension == "mp4":
    #     s = handle_mp4(file_path)
    elif extension == "pptx":
        s = handle_pptx(file_path)
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type")

    response = prompt_everyting(s)
    # quiz2 = generate_quiz(response["flash_cards"])
    # response["quiz"] += quiz2
    # random.shuffle(response["quiz"])

    # id = generate_id()
    # response["id"] = id
    # with open(os.getcwd() + "/database/" + id + ".json", "w") as f:
    #     json.dump(response, f)

    return response

# @app.post("/fetch_id")
# async def fetch_id(request: dict):
#     id = request.get("id")
#     try:
#         with open(os.getcwd() + "/database/" + id + ".json", "r") as f:
#             data = json.load(f)
#         return data
#     except FileNotFoundError:
#         raise HTTPException(status_code=404, detail="ID not found")

# @app.get("/recent")
# async def recent():
#     output = {"recent": []}
#     i = 0
#     for file in os.listdir(os.getcwd() + "/database"):
#         if i == 10:
#             break
#         if file.endswith(".json"):
#             with open(os.getcwd() + "/database/" + file, "r") as f:
#                 data = json.load(f)
#                 id = file.split(".json")[0]
#                 title = data["title"]
#                 output["recent"].append({"id": id, "title": title})
#         i += 1
#     return output

# @app.post("/export")
# async def export(request: dict):
#     selected = request.get("selected")
#     data = request.get("data")
    
#     if selected == 0:
#         export_summary(data["summary"], "Summary.docx")
#         return FileResponse("Summary.docx", filename="Summary.docx")
#     elif selected == 1:
#         export_flashcards(data["flash_cards"], "Flashcards.docx")
#         return FileResponse("Flashcards.docx", filename="Flashcards.docx")
#     else:
#         export_quiz(data["quiz"], "Quiz.docx")
#         return FileResponse("Quiz.docx", filename="Quiz.docx")

@app.post("/export")
async def export(request: Request):
    req = await request.json()
    selected = req.get("selected")
    data = req.get("data")
    
    filename = ""
    if selected == 0:
        filename = "Summary.docx"
        export_summary(data, filename)
    elif selected == 1:
        filename = "Flashcards.docx"
        export_flashcards(data, filename)
    else:
        filename = "Quiz.docx"
        export_quiz(data, filename)
    
    return FileResponse(path=filename, filename=filename, media_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=5000)